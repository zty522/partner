"""Stage report generation for long-running Partner projects.

The agent writes a compact Markdown report; Partner turns it into user-facing
PPTX/PDF artifacts. This keeps the LLM prompt light while still giving users a
readable "project meeting" output after enough work has accumulated.
"""

from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
from dataclasses import dataclass, asdict
from datetime import datetime, timedelta
from typing import Any
from zipfile import ZipFile


DEFAULT_EVERY_ROUNDS = 8
DEFAULT_MIN_INTERVAL_HOURS = 2
REQUIRED_SECTIONS = [
    "一句话结论",
    "背景与原始目标",
    "本阶段真正完成了什么",
    "关键证据与执行记录",
    "重要发现或判断变化",
    "失败、风险与不能声称的内容",
    "形成的新习惯/成长",
    "对用户有价值的产物",
    "下一步最小可验证动作",
]


@dataclass
class StageReport:
    project: str
    title: str
    created_at: str
    sections: list[dict[str, Any]]
    evidence_files: list[str]
    evidence_levels: list[str]
    source_markdown: str = ""
    validation: dict[str, Any] | None = None

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass
class ReportValidation:
    ok: bool
    errors: list[str]
    warnings: list[str]
    stats: dict[str, Any]

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


def _now() -> datetime:
    return datetime.now()


def _safe_name(name: str) -> str:
    name = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", name or "project")
    return name.strip("_") or "project"


def _state_path(workspace: str) -> str:
    state_dir = os.path.join(workspace, "state")
    os.makedirs(state_dir, exist_ok=True)
    return os.path.join(state_dir, "stage_report_state.json")


def _load_state(workspace: str) -> dict[str, Any]:
    try:
        with open(_state_path(workspace), "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def _save_state(workspace: str, state: dict[str, Any]) -> None:
    with open(_state_path(workspace), "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=2)


def _project_report_dir(workspace: str, project: str) -> str:
    from .projects.project_state import get_project_dir

    report_dir = os.path.join(get_project_dir(workspace, project), "reports")
    os.makedirs(report_dir, exist_ok=True)
    return report_dir


def _user_report_dir(workspace: str, project: str) -> str:
    report_dir = os.path.join(workspace, "state", "user", "reports", _safe_name(project))
    os.makedirs(report_dir, exist_ok=True)
    return report_dir


def _env_int(name: str, default: int) -> int:
    try:
        value = int(os.getenv(name, "").strip())
        return value if value > 0 else default
    except Exception:
        return default


def maybe_stage_report_objective(workspace: str, project: str, step: int) -> tuple[str, str]:
    """Return a report-writing objective when a project has enough accumulated work."""
    every_rounds = _env_int("PARTNER_STAGE_REPORT_EVERY_ROUNDS", DEFAULT_EVERY_ROUNDS)
    min_interval_hours = _env_int("PARTNER_STAGE_REPORT_MIN_INTERVAL_HOURS", DEFAULT_MIN_INTERVAL_HOURS)
    if step <= 0 or step % every_rounds != 0:
        return "", ""

    state = _load_state(workspace)
    project_state = state.get(project) if isinstance(state.get(project), dict) else {}
    last_at = str(project_state.get("last_generated_at") or "")
    if last_at:
        try:
            last_dt = datetime.fromisoformat(last_at)
            if _now() - last_dt < timedelta(hours=min_interval_hours):
                return "", ""
        except Exception:
            pass

    ts = _now().strftime("%Y%m%d_%H%M")
    path = os.path.join(_project_report_dir(workspace, project), f"stage_report_{ts}.md")
    try:
        from .projects.project_state import load_project_guardrail

        guardrail = load_project_guardrail(workspace, project)
    except Exception:
        guardrail = {}
    contract_clause = ""
    if guardrail:
        allowed = "；".join(
            [str(guardrail.get("current_mainline") or "").strip()]
            + [str(x).strip() for x in (guardrail.get("allowed_scope") or []) if str(x).strip()]
        ).strip("；")
        forbidden = "；".join(str(x).strip() for x in (guardrail.get("forbidden_scope") or []) if str(x).strip())
        criteria = "；".join(str(x).strip() for x in (guardrail.get("completion_criteria") or []) if str(x).strip())
        contract_clause = (
            "汇报必须对照用户原始任务合同。"
            f"允许/主线：{allowed or '按用户原始要求'}。"
            f"禁止/越界：{forbidden or '用户没有要求的新方向'}。"
            f"完成标准：{criteria or '直接满足用户点名交付物'}。"
            "本阶段真正完成了什么只能写合同内成果；越界内容只能放在风险或下一阶段建议里，不能写成本任务成果。"
        )
    objective = (
        "做一次高质量阶段汇报，不要继续扩展实验或大范围检索。基于 project_brief、state、exploration_log、"
        "长期记忆、growth_journal、habit_applications、runtime_cost 和最近产物，写一份面向用户/老师的 Markdown 组会式汇报。"
        f"{contract_clause}"
        "汇报必须具体、完整、诚实，不能只写几句摘要。建议 1600-2600 字。"
        "必须包含这些二级标题，并按顺序写："
        "1. 一句话结论；"
        "2. 背景与原始目标；"
        "3. 本阶段真正完成了什么；"
        "4. 关键证据与执行记录；"
        "5. 重要发现或判断变化；"
        "6. 失败、风险与不能声称的内容；"
        "7. 形成的新习惯/成长；"
        "8. 对用户有价值的产物；"
        "9. 下一步最小可验证动作。"
        "每个结论都要标注证据等级：REAL / LOCAL_EXECUTION / SIMULATION / INFERRED / BLOCKED。"
        "如果有代码执行，必须说明执行了什么、输出是什么、能证明什么、不能证明什么。"
        "如果生成了 PPT/PDF/showcase，要说明用户应该先看哪里。"
        "不要堆文件名、目录结构、字节数；不要问用户选方向；如果结果可能有泄露、幻觉或不可复现，要明确写出。"
    )
    return objective, path


def mark_stage_report_generated(workspace: str, project: str, markdown_path: str, outputs: dict[str, str]) -> None:
    state = _load_state(workspace)
    state[project] = {
        "last_generated_at": _now().isoformat(timespec="seconds"),
        "markdown": markdown_path,
        "json": outputs.get("json", ""),
        "pptx": outputs.get("pptx", ""),
        "pdf": outputs.get("pdf", ""),
        "pipeline": "structured_report_v1",
    }
    _save_state(workspace, state)


def _read_text(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception:
        return ""


def _is_placeholder_report(text: str) -> bool:
    raw = re.sub(r"\s+", " ", text or "").strip()
    if not raw:
        return True
    bad_patterns = [
        r"此处不重复全文",
        r"建议用户直接打开",
        r"请直接打开",
        r"报告已写入.*不重复",
        r"详见.*stage_report",
        r"见.*文件",
        r"不再重复",
    ]
    if len(raw) < 800:
        return True
    if len(raw) < 1600 and any(re.search(pattern, raw, re.I) for pattern in bad_patterns):
        return True
    # Long reports can legitimately mention generated files; do not reject
    # them unless the whole body is essentially a pointer.
    pointer_ratio = sum(1 for pattern in bad_patterns if re.search(pattern, raw, re.I))
    return len(raw) < 2400 and pointer_ratio >= 2


def _section_count(text: str) -> int:
    return len(re.findall(r"(?m)^##\s+", text or ""))


def _canonical_section(title: str) -> str:
    title = re.sub(r"^\s*\d+[.、]\s*", "", title or "").strip()
    title = re.sub(r"\s+", "", title)
    for required in REQUIRED_SECTIONS:
        key = re.sub(r"\s+", "", required)
        if key in title or title in key:
            return required
    return title or "摘要"


def _extract_evidence_files(text: str) -> list[str]:
    files = []
    for match in re.findall(r"[\w\u4e00-\u9fff./-]+\.(?:json|csv|md|py|txt|pdf|pptx|xlsx|docx|png|jpg|jpeg|gif)", text or ""):
        item = match.strip("`，,。；;()（）[]【】")
        if item and item not in files:
            files.append(item)
    return files[:60]


def _extract_evidence_levels(text: str) -> list[str]:
    levels = []
    for level in ("REAL", "LOCAL_EXECUTION", "SIMULATION", "INFERRED", "BLOCKED"):
        if re.search(rf"\b{level}\b", text or "") and level not in levels:
            levels.append(level)
    return levels


def stage_report_from_markdown(workspace: str, project: str, markdown_path: str) -> StageReport:
    markdown_path = _resolve_full_markdown_path(workspace, project, markdown_path)
    md = _read_text(markdown_path)
    title, parsed_sections = _parse_markdown(md)
    sections = []
    for raw_title, lines in parsed_sections:
        body = "\n".join(line for line in lines if line.strip()).strip()
        if not body:
            continue
        sections.append({
            "title": _canonical_section(raw_title),
            "body": body,
            "bullets": [line for line in lines if line.strip()],
        })
    report = StageReport(
        project=project,
        title=title or f"{project} 阶段汇报",
        created_at=_now().isoformat(timespec="seconds"),
        sections=sections,
        evidence_files=_extract_evidence_files(md),
        evidence_levels=_extract_evidence_levels(md),
        source_markdown=markdown_path,
    )
    validation = validate_stage_report(report)
    report.validation = validation.to_dict()
    if not validation.ok:
        raise RuntimeError(
            "stage report validation failed: "
            + "; ".join(validation.errors)
            + f" source={markdown_path}"
        )
    return report


def validate_stage_report(report: StageReport) -> ReportValidation:
    errors: list[str] = []
    warnings: list[str] = []
    section_titles = [_canonical_section(str(sec.get("title") or "")) for sec in report.sections]
    body = "\n\n".join(str(sec.get("title", "")) + "\n" + str(sec.get("body", "")) for sec in report.sections)
    plain_len = len(re.sub(r"\s+", "", body))
    if plain_len < 2200:
        errors.append(f"report body too short: {plain_len} chars")
    if _is_placeholder_report(body):
        errors.append("report body looks like a placeholder/pointer")
    missing = [sec for sec in REQUIRED_SECTIONS if sec not in section_titles]
    if missing:
        errors.append("missing required sections: " + ", ".join(missing[:5]))
    if len(report.evidence_levels) < 1:
        errors.append("missing evidence level markers")
    if len(report.evidence_files) < 2:
        warnings.append("few evidence files referenced")
    if not re.search(r"(风险|不能声称|失败|泄露|不可复现|BLOCKED|INFERRED)", body, re.I):
        errors.append("missing risk/limitation boundary")
    if not re.search(r"(下一步|最小|验证|动作|实验|审计|复现)", body):
        errors.append("missing executable next-step language")
    if re.search(r"(此处不重复全文|建议用户直接打开|报告已写入.*不重复)", body):
        errors.append("contains placeholder wording")
    stats = {
        "body_chars": plain_len,
        "sections": len(report.sections),
        "required_sections_present": len(REQUIRED_SECTIONS) - len(missing),
        "evidence_files": len(report.evidence_files),
        "evidence_levels": report.evidence_levels,
    }
    return ReportValidation(ok=not errors, errors=errors, warnings=warnings, stats=stats)


def render_stage_report_markdown(report: StageReport) -> str:
    lines = [
        f"# {report.title}",
        "",
        f"**项目：** {report.project}",
        f"**生成时间：** {report.created_at}",
        "",
    ]
    for section in report.sections:
        title = str(section.get("title") or "摘要").strip()
        body = str(section.get("body") or "").strip()
        lines.extend([f"## {title}", "", body, ""])
    if report.evidence_files:
        lines.extend(["## 证据索引", ""])
        for item in report.evidence_files[:30]:
            lines.append(f"- `{item}`")
        lines.append("")
    if report.validation:
        lines.extend(["## 报告质量检查", ""])
        stats = report.validation.get("stats", {}) if isinstance(report.validation, dict) else {}
        lines.append(f"- 正文字符数：{stats.get('body_chars', '')}")
        lines.append(f"- 必需章节覆盖：{stats.get('required_sections_present', '')}/{len(REQUIRED_SECTIONS)}")
        lines.append(f"- 证据等级：{', '.join(report.evidence_levels) if report.evidence_levels else 'EMPTY'}")
        warnings = report.validation.get("warnings", []) if isinstance(report.validation, dict) else []
        if warnings:
            lines.append(f"- 警告：{'；'.join(warnings)}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def _resolve_full_markdown_path(workspace: str, project: str, markdown_path: str) -> str:
    """Return a real report body, not a short pointer/placeholder.

    Agents sometimes write ARTIFACT_CONTENT as "the report is in X, not repeated".
    PPT/PDF must never be generated from that placeholder.  This resolver
    prefers the provided file only if it contains enough report body; otherwise
    it searches the project directory for a fuller stage_report markdown.
    """
    current = _read_text(markdown_path)
    if current and not _is_placeholder_report(current) and _section_count(current) >= 5:
        return markdown_path

    from .projects.project_state import get_project_dir

    candidates: list[tuple[int, str]] = []
    project_dir = get_project_dir(workspace, project)
    basename = os.path.basename(markdown_path or "")
    for root, _, files in os.walk(project_dir):
        for name in files:
            if not (name.startswith("stage_report_") and name.endswith(".md")):
                continue
            path = os.path.join(root, name)
            if os.path.abspath(path) == os.path.abspath(markdown_path):
                continue
            # Prefer same timestamp/name outside reports/, then largest body.
            score = 0
            if name == basename:
                score += 100_000
            if f"{os.sep}reports{os.sep}" not in path:
                score += 50_000
            text = _read_text(path)
            if _is_placeholder_report(text) or _section_count(text) < 5:
                continue
            score += len(text)
            candidates.append((score, path))
    if candidates:
        candidates.sort(reverse=True)
        resolved = candidates[0][1]
        # Replace the placeholder file with the full body so user/latest copies
        # and future publishes all point to a useful Markdown report.
        try:
            if os.path.abspath(resolved) != os.path.abspath(markdown_path):
                shutil.copy2(resolved, markdown_path)
        except Exception:
            pass
        return resolved
    if current:
        raise RuntimeError(
            f"stage report markdown is a placeholder or too short; refusing to publish: {markdown_path}"
        )
    raise RuntimeError(f"stage report markdown is empty or unreadable: {markdown_path}")


def _plain(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text or "")
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"[*_#>-]+", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _parse_markdown(md: str) -> tuple[str, list[tuple[str, list[str]]]]:
    lines = md.splitlines()
    title = "阶段汇报"
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_lines
        if current_title or current_lines:
            cleaned = [_plain(x) for x in current_lines if _plain(x)]
            sections.append((current_title or "摘要", cleaned[:14]))
        current_title = ""
        current_lines = []

    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if line.startswith("# "):
            if title == "阶段汇报":
                title = _plain(line[2:])
            continue
        if line.startswith("## "):
            flush()
            current_title = _plain(line[3:])
            continue
        if line.startswith("### "):
            current_lines.append(_plain(line[4:]))
            continue
        current_lines.append(line)
    flush()
    if not sections:
        sections = [("摘要", [_plain(x) for x in lines if _plain(x)][:8])]
    return title, sections[:12]


def _chunks(items: list[str], size: int = 4) -> list[list[str]]:
    return [items[i:i + size] for i in range(0, len(items), size)] or [[]]


def _add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=None):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return box


def _pptx_text_char_count(path: str) -> int:
    try:
        total = 0
        with ZipFile(path) as zf:
            for name in zf.namelist():
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    xml = zf.read(name).decode("utf-8", errors="ignore")
                    total += sum(len(x) for x in re.findall(r"<a:t>(.*?)</a:t>", xml))
        return total
    except Exception:
        return 0


def _assert_nonempty_artifact(path: str, kind: str, min_size: int, min_text_chars: int = 0) -> None:
    if not path or not os.path.exists(path):
        raise RuntimeError(f"{kind} artifact was not created: {path}")
    size = os.path.getsize(path)
    if size < min_size:
        raise RuntimeError(f"{kind} artifact is too small ({size} bytes): {path}")
    if kind.lower() == "pptx":
        text_chars = _pptx_text_char_count(path)
        if text_chars < min_text_chars:
            raise RuntimeError(f"PPTX appears blank ({text_chars} text chars): {path}")


def _generate_pptx_from_report(report: StageReport, output_path: str) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    title = report.title
    sections = [
        (str(section.get("title") or "摘要"), [str(x) for x in (section.get("bullets") or []) if str(x).strip()])
        for section in report.sections
    ]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    navy = (20, 38, 58)
    amber = (199, 119, 35)
    red = (172, 48, 45)
    gray = (90, 95, 102)

    slide = prs.slides.add_slide(blank)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 246, 240)
    _add_textbox(slide, Inches(0.7), Inches(0.7), Inches(11.8), Inches(0.8), report.project, 20, False, amber)
    _add_textbox(slide, Inches(0.7), Inches(1.55), Inches(11.8), Inches(1.2), title, 34, True, navy)
    _add_textbox(slide, Inches(0.75), Inches(3.05), Inches(11.2), Inches(0.7), "阶段性科研汇报：结论、证据等级、执行记录、风险边界与成长习惯", 22, False, gray)
    _add_textbox(slide, Inches(0.75), Inches(5.9), Inches(11), Inches(0.4), datetime.now().strftime("%Y-%m-%d %H:%M"), 14, False, gray)

    for idx, (section_title, lines) in enumerate(sections, start=1):
        for part_idx, group in enumerate(_chunks(lines, 4), start=1):
            slide = prs.slides.add_slide(blank)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            title_text = f"{idx:02d}. {section_title}" if part_idx == 1 else f"{idx:02d}. {section_title}（续）"
            color = red if re.search(r"风险|审计|泄露|失败|边界", section_title) else navy
            _add_textbox(slide, Inches(0.65), Inches(0.45), Inches(11.8), Inches(0.55), title_text, 25, True, color)
            line = slide.shapes.add_shape(1, Inches(0.65), Inches(1.15), Inches(1.25), Inches(0.06))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*amber)
            line.line.color.rgb = RGBColor(*amber)

            top = 1.55
            for item in group:
                shape = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.1), Inches(0.7))
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item[:300]
                p.level = 0
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(19)
                p.font.color.rgb = RGBColor(36, 42, 50)
                top += 1.08
            _add_textbox(slide, Inches(0.65), Inches(6.95), Inches(11.8), Inches(0.25), "Partner 自动阶段汇报", 10, False, gray)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    _assert_nonempty_artifact(output_path, "pptx", min_size=10000, min_text_chars=1200)


def _generate_pptx(markdown_path: str, project: str, output_path: str) -> None:
    report = StageReport(
        project=project,
        title=_parse_markdown(_read_text(markdown_path))[0],
        created_at=_now().isoformat(timespec="seconds"),
        sections=[
            {"title": title, "body": "\n".join(lines), "bullets": lines}
            for title, lines in _parse_markdown(_read_text(markdown_path))[1]
        ],
        evidence_files=_extract_evidence_files(_read_text(markdown_path)),
        evidence_levels=_extract_evidence_levels(_read_text(markdown_path)),
        source_markdown=markdown_path,
    )
    _generate_pptx_from_report(report, output_path)


def _generate_pdf_from_report(report: StageReport, output_path: str) -> None:
    import os as _pdf_os
    _pdf_os.environ.setdefault("OPENBLAS_NUM_THREADS", "1")
    _pdf_os.environ.setdefault("OMP_NUM_THREADS", "1")
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer
    from reportlab.lib.utils import ImageReader

    # Resolve source markdown directory for relative image paths
    _md_dir = os.path.dirname(report.source_markdown) if report.source_markdown else ""

    def pick_cjk_font() -> str:
        candidates = [
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc",
            "/usr/share/fonts/truetype/arphic/uming.ttc",
            "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
            "/mnt/c/Windows/Fonts/msyh.ttc",
            "/mnt/c/Windows/Fonts/simsun.ttc",
        ]
        for candidate in candidates:
            if not os.path.exists(candidate):
                continue
            try:
                pdfmetrics.registerFont(TTFont("PartnerCJK", candidate))
                return "PartnerCJK"
            except Exception:
                continue
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            return "STSong-Light"
        except Exception:
            return "Helvetica"

    font_name = pick_cjk_font()

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="CNTitle", fontName=font_name, fontSize=22, leading=28, textColor=colors.HexColor("#14263a"), spaceAfter=16))
    styles.add(ParagraphStyle(name="CNHeading", fontName=font_name, fontSize=15, leading=20, textColor=colors.HexColor("#ac302d"), spaceBefore=14, spaceAfter=8))
    styles.add(ParagraphStyle(name="CNBody", fontName=font_name, fontSize=10.5, leading=16, textColor=colors.HexColor("#242a32"), leftIndent=8))
    styles.add(ParagraphStyle(name="CNMeta", fontName=font_name, fontSize=9, leading=13, textColor=colors.HexColor("#666666"), spaceAfter=10))
    story = [
        Paragraph(report.project, styles["CNMeta"]),
        Paragraph(report.title, styles["CNTitle"]),
        Paragraph(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}", styles["CNMeta"]),
    ]

    # Build sections from report data — include FULL body text, not just bullets
    for idx, section in enumerate(report.sections, start=1):
        section_title = str(section.get("title") or "摘要")
        body = str(section.get("body") or "")
        bullets = [str(x) for x in (section.get("bullets") or []) if str(x).strip()]
        story.append(Paragraph(f"{idx:02d}. {section_title}", styles["CNHeading"]))
        # Show body first (the full paragraph text)
        if body:
            for para in body.split("\n\n"):
                para = para.strip()
                if para and para not in bullets:
                    story.append(Paragraph(para, styles["CNBody"]))
                    story.append(Spacer(1, 0.08 * cm))
        # Then list bullet points
        for item in bullets or ["待补充。"]:
            story.append(Paragraph("• " + item, styles["CNBody"]))
            story.append(Spacer(1, 0.12 * cm))
        if idx in {3, 6}:
            story.append(PageBreak())

    # Embed images from evidence_files
    _max_img_w = 14 * cm  # max image width (fit A4 with margins)
    _max_img_h = 18 * cm
    for fpath in report.evidence_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        # Strip markdown image syntax: ![alt](path) -> path
        md_img = re.match(r"\[.*?\]\((.+?)\)", fpath)
        if md_img:
            fpath = md_img.group(1)
        ext = os.path.splitext(fpath)[1].lower()
        if ext not in (".png", ".jpg", ".jpeg", ".gif"):
            continue

        # Resolve image path: try as-is, then relative to markdown dir
        resolved = None
        for candidate in [fpath, os.path.join(_md_dir, fpath)]:
            if os.path.exists(candidate):
                resolved = candidate
                break
        if not resolved or not os.path.isfile(resolved):
            continue
        # Skip zero-byte images
        if os.path.getsize(resolved) < 50:
            continue
        try:
            img = Image(resolved)
            # Scale to fit page width while preserving aspect ratio
            iw, ih = img.imageWidth, img.imageHeight
            if iw > _max_img_w:
                ratio = _max_img_w / iw
                img.drawWidth = _max_img_w
                img.drawHeight = ih * ratio
            if img.drawHeight > _max_img_h:
                ratio = _max_img_h / img.drawHeight
                img.drawHeight = _max_img_h
                img.drawWidth = img.drawWidth * ratio
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)
            story.append(Spacer(1, 0.3 * cm))
        except Exception:
            continue

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    SimpleDocTemplate(output_path, pagesize=A4, rightMargin=1.5 * cm, leftMargin=1.5 * cm, topMargin=1.3 * cm, bottomMargin=1.3 * cm).build(story)
    _assert_nonempty_artifact(output_path, "pdf", min_size=8000)


def _generate_pdf_from_markdown(markdown_path: str, project: str, output_path: str) -> None:
    md = _read_text(markdown_path)
    title, sections = _parse_markdown(md)
    report = StageReport(
        project=project,
        title=title,
        created_at=_now().isoformat(timespec="seconds"),
        sections=[
            {"title": section_title, "body": "\n".join(lines), "bullets": lines}
            for section_title, lines in sections
        ],
        evidence_files=_extract_evidence_files(md),
        evidence_levels=_extract_evidence_levels(md),
        source_markdown=markdown_path,
    )
    _generate_pdf_from_report(report, output_path)


def _try_convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        return False
    outdir = os.path.dirname(pdf_path)
    try:
        result = subprocess.run(
            [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, pptx_path],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=90,
            check=False,
        )
        converted = os.path.join(outdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        if result.returncode == 0 and os.path.exists(converted):
            if os.path.abspath(converted) != os.path.abspath(pdf_path):
                shutil.move(converted, pdf_path)
            return True
    except Exception:
        return False
    return False


def publish_stage_report(workspace: str, project: str, markdown_path: str) -> dict[str, str]:
    """Create PPTX/PDF copies for a Markdown stage report and expose them to user/."""
    if not markdown_path or not os.path.exists(markdown_path):
        return {}
    report = stage_report_from_markdown(workspace, project, markdown_path)
    report_dir = _project_report_dir(workspace, project)
    ts = _now().strftime("%Y%m%d_%H%M")
    base = f"{_safe_name(project)}_stage_report_{ts}"
    canonical_md_path = os.path.join(report_dir, base + ".md")
    report_json_path = os.path.join(report_dir, base + ".json")
    pptx_path = os.path.join(report_dir, base + ".pptx")
    pdf_path = os.path.join(report_dir, base + ".pdf")

    os.makedirs(report_dir, exist_ok=True)
    with open(canonical_md_path, "w", encoding="utf-8") as f:
        f.write(render_stage_report_markdown(report))
    with open(report_json_path, "w", encoding="utf-8") as f:
        json.dump(report.to_dict(), f, ensure_ascii=False, indent=2)

    _generate_pptx_from_report(report, pptx_path)
    # Generate PDF directly from Markdown. LibreOffice conversion can produce
    # visually blank PDFs on minimal Linux images when CJK fonts are missing.
    _generate_pdf_from_report(report, pdf_path)

    user_dir = _user_report_dir(workspace, project)
    user_md = os.path.join(user_dir, "latest_stage_report.md")
    user_json = os.path.join(user_dir, "latest_stage_report.json")
    user_pptx = os.path.join(user_dir, os.path.basename(pptx_path))
    user_pdf = os.path.join(user_dir, os.path.basename(pdf_path))
    shutil.copy2(canonical_md_path, user_md)
    shutil.copy2(report_json_path, user_json)
    shutil.copy2(pptx_path, user_pptx)
    shutil.copy2(pdf_path, user_pdf)
    readme_path = os.path.join(user_dir, "README.md")
    with open(readme_path, "w", encoding="utf-8") as f:
        f.write(
            f"# {project} 阶段汇报\n\n"
            "- `latest_stage_report.md`：最新文字版汇报\n"
            "- `latest_stage_report.json`：结构化汇报对象和质量检查\n"
            f"- `{os.path.basename(user_pptx)}`：最新 PPT 汇报\n"
            f"- `{os.path.basename(user_pdf)}`：最新 PDF 汇报\n"
        )
    outputs = {"markdown": user_md, "json": user_json, "pptx": user_pptx, "pdf": user_pdf}
    mark_stage_report_generated(workspace, project, canonical_md_path, outputs)
    try:
        from .showcase import build_showcase

        showcase_dir = build_showcase(workspace, project)
        outputs["showcase"] = str(showcase_dir)
    except Exception:
        pass
    return outputs
